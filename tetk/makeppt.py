from pathlib import Path
import pandas as pd
import copy
import pptx
from pptx.util import Cm, Pt
from itertools import zip_longest
import PIL
import math

import commons as cf
from reader import FileManager


class PptManager:
    def __init__(self, cfg: dict, filepath_template: Path) -> None:
        self.cfg = cfg
        assert filepath_template.is_file(), f"invalid {filepath_template=}"
        self.root = pptx.Presentation(filepath_template)

    @staticmethod
    def split_list(list_in, no_of_chunks, fillvalue="none"):
        return zip_longest(*[iter(list_in)] * no_of_chunks, fillvalue=fillvalue)

    def crop_image(
        self, imgpath: Path, orignal_size=None, crop_coords=(0, 0, 0, 0)
    ) -> bool:
        if not imgpath.is_file():
            raise Exception(f"invalid image path: {imgpath}")
        img = PIL.Image.open(imgpath)
        if orignal_size is not None:
            if (img.size[0] == orignal_size[0]) and (img.size[1] == orignal_size[1]):
                pass
            else:
                self.cfg.log.debug(f"Skip crop. Non-original img: {imgpath.name}")
                return False

        cropped_img = img.crop(crop_coords)
        cropped_img.save(imgpath)
        self.cfg.log.info(f"Cropped and overridden image: {imgpath.name}")
        return True

    def insert_images_wafersPerLot(
        self, name: str, images: list, stepsize_x=3, stepsize_y=3, fixed_width=3
    ):
        log = self.cfg.log
        root = self.root
        slide = root.slides.add_slide(root.slide_layouts[3])
        cols = math.ceil(math.sqrt(len(images)))
        n = len(images)
        images = list(self.split_list(images, cols))
        legend_text = "Legend:\n"
        coor_x_start = self.cfg["coor_x_origin"]
        coor_x = copy.copy(coor_x_start)
        coor_y_start = self.cfg["coor_y_origin"]
        coor_y = copy.copy(coor_y_start)

        title_txtbox = slide.shapes.add_textbox(
            Cm(coor_x_start),
            Cm(coor_y_start + self.cfg["wafer_id_text_y_offset"]),
            width=Cm(self.cfg["wafer_id_textbox_width_cm"]),
            height=Cm(self.cfg["wafer_id_textbox_height_cm"]),
        )
        title = title_txtbox.text_frame.paragraphs[0]
        title.text = f"{name}: {n} wafer images "
        title.font.size = Pt(self.cfg["wafer_id_text_size"])

        image_counter = 0
        for row in images:
            for imgpath in row:
                if not imgpath or imgpath == "none":
                    continue
                self.crop_image(
                    imgpath,
                    orignal_size=self.cfg["original_image_size"],
                    crop_coords=self.cfg["crop_coords"],
                )
                slide.shapes.add_picture(
                    str(imgpath.resolve()),
                    left=Cm(coor_x),
                    top=Cm(coor_y),
                    width=Cm(fixed_width),
                )
                wafer_id = imgpath.name.split("-")[-2]
                try:
                    wafer_id = int(wafer_id)
                except Exception:
                    wafer_id = 0
                legend_text += f" {wafer_id:02d}"
                log.debug(f"added {imgpath.name=}")
                coor_x += stepsize_x
                image_counter += 1
            legend_text += "\n"

            coor_y += stepsize_y
            coor_x = copy.copy(coor_x_start)

        legend_txtbox = slide.shapes.add_textbox(
            Cm(coor_x_start + self.cfg["wafer_id_legend_x_offset_cm"]),
            Cm(coor_y_start + self.cfg["wafer_id_legend_y_offset_cm"]),
            width=Cm(3),
            height=Cm(3),
        )
        legend_txtbox.text_frame.word_wrap = False
        legend = legend_txtbox.text_frame.paragraphs[0]
        legend.text = legend_text
        legend.font.size = Pt(self.cfg["wafer_id_text_size"])
        legend.font.name = "Courier"

    def insert_images_singleWaferCompare(
        self, df: pd.DataFrame, cols=4, stepsize_x=6, stepsize_y=6, fixed_width=6
    ):
        root = self.root
        foldernames = list(df.columns)
        col_count = 0
        legend_headerTop = []
        legend_headerBtm = []
        legend_firstRow = []
        legend_secondRow = []

        coor_x_start = self.cfg["coor_x_origin"]
        coor_x = copy.copy(coor_x_start)
        coor_y_start = self.cfg["coor_y_origin"]
        coor_y = copy.copy(coor_y_start)

        for i, (ind, row) in enumerate(df.iterrows()):
            if col_count == 0:
                slide = root.slides.add_slide(root.slide_layouts[3])
                coor_x = copy.copy(coor_x_start)
                coor_y = copy.copy(coor_y_start)

            for j, foldername in enumerate(foldernames):
                imgpath = row[foldername]
                try:
                    self.crop_image(
                        imgpath,
                        orignal_size=self.cfg["original_image_size"],
                        crop_coords=self.cfg["crop_coords"],
                    )
                    slide.shapes.add_picture(
                        str(imgpath.resolve()),
                        left=Cm(coor_x),
                        top=Cm(coor_y),
                        width=Cm(fixed_width),
                    )
                except Exception as e:
                    self.log.error(f"insert_images_singleWaferCompare(); {e=}")
                if j == 0:
                    legend_headerTop.append("::".join(ind))
                    legend_firstRow.append(foldername)
                    coor_y += stepsize_y
                else:
                    legend_headerBtm.append("::".join(ind))
                    legend_secondRow.append(foldername)
                    coor_x += stepsize_x
                    coor_y -= stepsize_y

            self.cfg.log.info(
                f"singleRun{i:03d}: {ind}, {foldernames[0]}={bool(row[foldernames[0]])} {foldernames[1]}={bool(row[foldernames[1]])}"
            )

            col_count += 1
            if col_count == 4 or i == len(df) - 1:
                col_count = 0
                legend_txtbox = slide.shapes.add_textbox(
                    Cm(
                        coor_x_start
                        + self.cfg["wafer_id_single_txtbox_loc_offset_x_cm"]
                    ),
                    Cm(
                        coor_y_start
                        + self.cfg["wafer_id_single_txtbox_loc_offset_y_cm"]
                    ),
                    width=Cm(3),
                    height=Cm(3),
                )
                legend = legend_txtbox.text_frame.paragraphs[0]
                legend_headerTop = [f"{v:^13}" for v in legend_headerTop]
                legend_firstRow = [f"{v:^13}" for v in legend_firstRow]
                legend_secondRow = [f"{v:^13}" for v in legend_secondRow]
                legend.text = f"Legend:\n{' '.join(legend_headerTop)}\n{' '.join(legend_firstRow)}\n{' '.join(legend_secondRow)}"
                legend.font.size = Pt(self.cfg["wafer_id_single_txtbox_text_size_pt"])
                legend.font.name = "Courier"
                legend_txtbox.text_frame.word_wrap = False
                legend_headerTop = []
                legend_headerBtm = []
                legend_firstRow = []
                legend_secondRow = []

    def save_ppt(self):
        outname = f"output-{cf.get_time()}.pptx"
        outpath = self.cfg.user_folders["outf01"] / outname
        self.root.save(outpath)
        cf.output(self.cfg.log, outpath)


def makeppt(cfg):
    fmgr = FileManager(cfg, cfg.user_folders["inpf01"])

    ppt = PptManager(cfg, cfg.user_files["ppt_template1"])

    # data = fmgr.construct_lotPerPage()
    # for i, (wafer_id, wafer_imgpaths) in enumerate(data.items()):
    #     cfg.log.info(f"{wafer_id=}; {len(wafer_imgpaths)=}")
    #     ppt.insert_images_wafersPerLot(wafer_id, wafer_imgpaths)

    data = fmgr.construct_singleWafers()
    ppt.insert_images_singleWaferCompare(
        df=data,
        stepsize_x=cfg["wafer_id_single_size_width_cm"],
        stepsize_y=cfg["wafer_id_single_size_width_cm"],
        fixed_width=cfg["wafer_id_single_size_width_cm"],
    )

    ppt.save_ppt()


if __name__ == "__main__":
    from main import cfg

    makeppt(cfg)
